"""Generate a concise Chinese presentation for the Manshu OCR/NER workflow.

The deck is intentionally curated rather than a one-to-one markdown dump: it
keeps the narrative useful for a short group presentation.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
OUT = DOCS / "MANSHU_OCR_NER_RESEARCH_WORKFLOW_PRESENTATION.pptx"
SOURCE_MD = DOCS / "MANSHU_OCR_NER_RESEARCH_WORKFLOW_CASE_STUDY.md"

FONT = "Microsoft YaHei"
BG = RGBColor(247, 244, 236)
INK = RGBColor(35, 40, 45)
MUTED = RGBColor(94, 102, 110)
ACCENT = RGBColor(20, 119, 111)
ACCENT_2 = RGBColor(201, 116, 52)
ACCENT_3 = RGBColor(44, 82, 130)
LIGHT = RGBColor(255, 253, 248)
LINE = RGBColor(217, 209, 195)


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor = LINE, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(slide, x, y, w, h, text="", font_size=24, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, kicker: str | None = None, page_no: int | None = None) -> None:
    if kicker:
        add_textbox(slide, 0.7, 0.42, 5.5, 0.35, kicker, 10, ACCENT_2, True)
    add_textbox(slide, 0.7, 0.72, 9.7, 0.78, title, 26, INK, True)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.58), Inches(1.0), Inches(0.08))
    set_fill(bar, ACCENT)
    bar.line.fill.background()
    if page_no is not None:
        add_textbox(slide, 12.15, 6.95, 0.55, 0.2, f"{page_no:02d}", 8, MUTED, False, PP_ALIGN.RIGHT)


def add_footer(slide) -> None:
    add_textbox(slide, 0.7, 6.95, 8.5, 0.22, "AI 辅助历史研究流程复盘：《満洲紳士録》PDF → OCR → NER", 8, MUTED)


def add_bullets(slide, x, y, w, h, items, font_size=15, color=INK, gap=0.12) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap * 28)
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.line_spacing = 1.12
        p.bullet = True


def add_card(slide, x, y, w, h, title, body, accent=ACCENT) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(card, LIGHT)
    set_line(card)
    stripe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.18), Inches(y + 0.2), Inches(0.08), Inches(h - 0.4))
    set_fill(stripe, accent)
    stripe.line.fill.background()
    add_textbox(slide, x + 0.38, y + 0.22, w - 0.55, 0.35, title, 14, INK, True)
    add_textbox(slide, x + 0.38, y + 0.72, w - 0.55, h - 0.9, body, 11, MUTED)


def add_big_number(slide, x, y, number, label, color=ACCENT) -> None:
    add_textbox(slide, x, y, 1.7, 0.6, number, 32, color, True, PP_ALIGN.CENTER)
    add_textbox(slide, x - 0.2, y + 0.68, 2.1, 0.45, label, 10, MUTED, False, PP_ALIGN.CENTER)


def add_table(slide, x, y, w, h, headers, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        set_fill(cell, ACCENT_3)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            set_fill(cell, LIGHT if r_idx % 2 else RGBColor(239, 236, 226))
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(10)
                p.font.color.rgb = INK
                p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
    return table_shape


def add_flow(slide, labels, y=3.0) -> None:
    x = 0.9
    box_w = 1.75
    gap = 0.36
    for idx, (label, accent) in enumerate(labels):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(0.68))
        set_fill(shape, LIGHT)
        set_line(shape, accent, 1.5)
        add_textbox(slide, x + 0.1, y + 0.2, box_w - 0.2, 0.2, label, 12, INK, True, PP_ALIGN.CENTER)
        if idx < len(labels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + box_w + 0.08), Inches(y + 0.22), Inches(0.22), Inches(0.24))
            set_fill(arrow, accent)
            arrow.line.fill.background()
        x += box_w + gap


def apply_background(slide) -> None:
    set_fill(slide.background, BG)
    blob = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(10.25), Inches(-0.5), Inches(3.0), Inches(2.4))
    blob.line.color.rgb = RGBColor(229, 221, 204)
    blob.line.width = Pt(2)
    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.55), Inches(0.28), Inches(0.92), Inches(0.12))
    set_fill(accent, ACCENT_2)
    accent.line.fill.background()


def add_slide(prs, title, kicker=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    add_title(slide, title, kicker, len(prs.slides))
    add_footer(slide)
    return slide


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    add_textbox(slide, 0.82, 0.72, 2.7, 0.28, "方法复盘 / 演示稿", 11, ACCENT_2, True)
    add_textbox(slide, 0.82, 1.35, 9.6, 1.4, "从版面切分到 NER", 40, INK, True)
    add_textbox(slide, 0.86, 2.48, 9.8, 0.55, "AI 辅助《満洲紳士録》人物资料处理流程", 21, MUTED)
    add_flow(
        slide,
        [("PDF", ACCENT_3), ("版面切分", ACCENT), ("ndlocr-lite", ACCENT_2), ("清洗分条", ACCENT), ("NER", ACCENT_3), ("研究数据", ACCENT_2)],
        4.18,
    )
    add_textbox(slide, 0.82, 6.65, 8.8, 0.25, "基于工作区处理记录与 docs 案例文档整理", 9, MUTED)

    # 2
    slide = add_slide(prs, "研究目标：把史料页变成可分析数据", "01 / 目标")
    add_bullets(
        slide,
        0.9,
        2.05,
        5.5,
        2.8,
        [
            "对象：工作区中已有的 manshu 史料 PDF。",
            "目标数据：人物条目、姓名、身份履历、地址、组织机构等结构化信息。",
            "原则：版面切分、OCR 后清洗、NER 尽可能高质量；OCR 主体使用本地 ndlocr-lite 控制成本。",
            "约束：不进行人工标注；大模型 API 只允许极小范围测试。",
        ],
    )
    add_card(slide, 7.0, 2.1, 4.8, 1.25, "核心判断", "NER 质量不是单靠 prompt 决定；上游版面与文本边界会直接决定下游能否“看见”信息。", ACCENT_2)
    add_card(slide, 7.0, 3.75, 4.8, 1.25, "流程心法", "先把页面结构恢复正确，再让 OCR 和 NER 各自做它们擅长的事。", ACCENT)

    # 3
    slide = add_slide(prs, "为什么这批材料难处理", "02 / 难点")
    add_card(slide, 0.9, 2.0, 3.55, 1.4, "版面", "双页扫描、竖排多栏、栏间距窄，错误切分会直接打乱阅读顺序。", ACCENT_3)
    add_card(slide, 4.9, 2.0, 3.55, 1.4, "文字", "日文旧字体、人名地名与官职密集，OCR 容错空间很小。", ACCENT)
    add_card(slide, 8.9, 2.0, 3.55, 1.4, "语义", "人物条目边界不总是显式，身份、组织、地址常交织在同一句中。", ACCENT_2)
    add_textbox(slide, 1.0, 4.25, 10.7, 0.8, "一个关键问题：如果切分阶段把竖排栏切碎或把左右页混在一起，后面的 OCR 和 NER 会“认真地处理错误输入”。", 19, INK, True, PP_ALIGN.CENTER)

    # 4
    slide = add_slide(prs, "前三轮尝试暴露出的问题", "03 / 早期流程")
    add_bullets(
        slide,
        0.9,
        2.0,
        5.5,
        3.4,
        [
            "不同切分策略产生的 OCR 文本顺序差异很大。",
            "部分 NER prompt 虽然字段很多，但抽取时漏掉人物、履历和组织信息。",
            "前 5 个条目仍有错漏，说明问题不是单点 bug，而是流程链条上多环节共同放大。",
            "需要从版面、清洗、分条、prompt、输出校验一起重审。",
        ],
    )
    add_card(slide, 7.0, 2.0, 4.7, 1.3, "发现", "失败样例最有价值：它们告诉我们信息到底是在 OCR 前丢了、清洗时丢了，还是 NER 时没抽出来。", ACCENT_2)
    add_card(slide, 7.0, 3.65, 4.7, 1.3, "转向", "不再只优化 prompt，而是把“结构恢复”作为 NER 的前置条件。", ACCENT)

    # 5
    slide = add_slide(prs, "关键诊断：先修版面，再谈语义", "04 / 诊断")
    add_bullets(
        slide,
        0.9,
        2.0,
        5.6,
        3.5,
        [
            "横向条带式切分容易破坏竖排栏，导致人物条目跨栏错序。",
            "双页扫描需要先拆成左右半页，再按半页做 OCR 与重构。",
            "NER 漏信息时，要回溯检查：原图是否有、OCR XML 是否有、清洗文本是否保留。",
            "以页为单位做小范围闭环，比盲目全书重跑更可靠。",
        ],
    )
    add_flow(slide, [("原图", ACCENT_3), ("半页", ACCENT), ("OCR XML", ACCENT_2), ("重构文本", ACCENT), ("分条", ACCENT_3), ("NER JSONL", ACCENT_2)], 4.95)

    # 6
    slide = add_slide(prs, "优化后的版面与 OCR 策略", "05 / 版面 + OCR")
    add_card(slide, 0.9, 2.0, 3.7, 1.45, "1. 双页拆分", "按左右半页恢复书页单元，避免左右页混读。", ACCENT_3)
    add_card(slide, 4.85, 2.0, 3.7, 1.45, "2. 本地 OCR", "使用 ndlocr-lite 读取日文印刷体，输出 XML，保留坐标与行块。", ACCENT)
    add_card(slide, 8.8, 2.0, 3.7, 1.45, "3. XML 重构", "用坐标与文本特征恢复竖排阅读顺序，而不是只拼纯文本。", ACCENT_2)
    add_textbox(slide, 1.05, 4.25, 10.9, 0.9, "这一步的目标不是“让 OCR 完美”，而是尽量让下游拿到顺序正确、边界清晰、信息不被提前删掉的文本。", 19, INK, True, PP_ALIGN.CENTER)

    # 7
    slide = add_slide(prs, "清洗与分条：避免把信息洗掉", "06 / 清洗")
    add_bullets(
        slide,
        0.9,
        2.0,
        5.7,
        3.5,
        [
            "保留标题行、坐标来源、页码和半页来源，方便追溯。",
            "清洗只处理明显噪声，不删除看似“多余”的官职、地址、组织信息。",
            "人物分条时允许前置官职、别名、组织名进入上下文。",
            "输出 review 标记：对疑似断裂、低置信、字段冲突条目留痕。",
        ],
    )
    add_card(slide, 7.1, 2.35, 4.55, 1.55, "清洗边界", "历史资料里“脏”和“有意义的异体表达”常常长得很像，所以清洗要保守，校验要主动。", ACCENT_2)

    # 8
    slide = add_slide(prs, "NER prompt 的新要求", "07 / NER")
    add_card(slide, 0.9, 2.05, 3.55, 1.25, "上下文", "以小 chunk 处理，但保留条目边界与页内来源，减少跨人物污染。", ACCENT_3)
    add_card(slide, 4.9, 2.05, 3.55, 1.25, "证据", "字段尽量带 evidence，方便知道抽取依据来自哪段原文。", ACCENT)
    add_card(slide, 8.9, 2.05, 3.55, 1.25, "审阅", "对不确定字段进入 review，而不是静默丢弃。", ACCENT_2)
    add_bullets(
        slide,
        1.0,
        4.05,
        10.7,
        1.45,
        [
            "重点字段：姓名、身份类别、组织机构、职务履历、地址、籍贯、关系、时间线。",
            "输出目标：既能机器读取，也能回到原文核查。",
        ],
        16,
    )

    # 9
    slide = add_slide(prs, "211-215 页小范围闭环验证", "08 / 验证")
    add_table(
        slide,
        1.0,
        2.0,
        8.7,
        2.0,
        ["页码", "分条数", "NER 人物数", "状态"],
        [
            ["211", "19", "19", "通过"],
            ["212", "19", "19", "通过"],
            ["213", "21", "21", "通过"],
            ["214", "21", "21", "通过"],
            ["215", "21", "21", "通过"],
        ],
        [2.0, 2.0, 2.2, 2.0],
    )
    add_big_number(slide, 10.25, 2.15, "101", "小范围条目 / 人物闭环", ACCENT)
    add_textbox(slide, 1.05, 5.05, 10.8, 0.55, "意义：先用 5 页验证流程链条，再扩展到全书，避免把系统性错误复制到全量结果中。", 17, INK, True, PP_ALIGN.CENTER)

    # 10
    slide = add_slide(prs, "本地 OCR 模型横向试验", "09 / OCR 对比")
    add_table(
        slide,
        1.0,
        2.05,
        9.6,
        2.05,
        ["模型", "第 211 页识别/分条表现", "适用判断"],
        [
            ["ndlocr-lite", "19 / 19", "当前主流程"],
            ["PaddleOCR", "11 / 19", "可作为补充参考"],
            ["RapidOCR", "8 / 19", "漏识较多"],
            ["EasyOCR", "0 / 19", "不适合该材料"],
        ],
        [2.4, 3.0, 4.2],
    )
    add_card(slide, 1.05, 4.75, 9.7, 0.9, "说明", "该轮只评价 OCR 结果，不使用大模型 API 做 NER，以避免把语义修复能力误判为 OCR 能力。", ACCENT_2)

    # 11
    slide = add_slide(prs, "下一步全书级流程", "10 / 落地流程")
    add_flow(
        slide,
        [
            ("PDF", ACCENT_3),
            ("页图导出", ACCENT),
            ("左右拆页", ACCENT_2),
            ("ndlocr-lite", ACCENT),
            ("XML 重构", ACCENT_3),
            ("清洗分条", ACCENT_2),
        ],
        2.35,
    )
    add_flow(
        slide,
        [
            ("抽样检查", ACCENT_3),
            ("Qwen NER", ACCENT),
            ("JSONL 合并", ACCENT_2),
            ("质量报告", ACCENT),
            ("人工复核点", ACCENT_3),
            ("研究数据", ACCENT_2),
        ],
        4.25,
    )
    add_textbox(slide, 1.0, 5.85, 10.7, 0.5, "全书处理前，优先固定输入输出目录、日志、错误重跑机制和质量统计表。", 16, MUTED, False, PP_ALIGN.CENTER)

    # 12
    slide = add_slide(prs, "方法论总结", "11 / 总结")
    add_bullets(
        slide,
        1.0,
        2.0,
        10.8,
        3.4,
        [
            "不要把 OCR、清洗、NER 看成孤立步骤；历史资料处理更像一条证据链。",
            "prompt 优化有效，但前提是上游文本边界和顺序没有被破坏。",
            "无人工标注条件下，小样本闭环、可追溯证据和 review 标记尤其重要。",
            "最适合全书扩展的方案，是成本可控、错误可定位、结果可复核的方案。",
        ],
        17,
    )
    add_textbox(slide, 1.1, 6.15, 10.6, 0.45, "从“让模型帮我识别”到“设计一条可审计的研究流水线”。", 19, ACCENT, True, PP_ALIGN.CENTER)

    return prs


def main() -> None:
    if not SOURCE_MD.exists():
        raise FileNotFoundError(f"Source markdown not found: {SOURCE_MD}")
    DOCS.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUT)
    print(f"Generated: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
